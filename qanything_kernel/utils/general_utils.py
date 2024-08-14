from sanic.request import Request
from sanic.exceptions import BadRequest
from qanything_kernel.utils.custom_log import debug_logger, embed_logger, rerank_logger
from qanything_kernel.configs.model_config import (KB_SUFFIX, UPLOAD_ROOT_PATH, LOCAL_EMBED_PATH, LOCAL_RERANK_PATH)
from transformers import AutoTokenizer
import pandas as pd
import inspect
import traceback
from urllib.parse import urlparse
import time
import os
import logging
import re
import requests
import aiohttp
from functools import wraps
import tiktoken
from openpyxl.utils import get_column_letter
from openpyxl import load_workbook
import numpy as np
from datetime import datetime, timedelta
from bs4 import BeautifulSoup
from urllib.parse import urljoin
import html2text
import os
import csv
import docx2txt
import fitz  # PyMuPDF
import openpyxl
from pptx import Presentation
import email
import chardet

__all__ = ['isURL', 'get_time', 'get_time_async', 'format_source_documents', 'safe_get', 'truncate_filename',
           'shorten_data', 'read_files_with_extensions', 'validate_user_id', 'get_invalid_user_id_msg', 'num_tokens',
           'clear_string', 'simplify_filename', 'string_bytes_length', 'correct_kb_id', 'clear_kb_id',
           'clear_string_is_equal', 'export_qalogs_to_excel', 'deduplicate_documents', 'fast_estimate_file_char_count',
           'check_user_id_and_user_info', 'get_table_infos', 'format_time_record', 'get_time_range',
           'html_to_markdown', "num_tokens_embed", "num_tokens_rerank", "get_all_subpages"]


def get_invalid_user_id_msg(user_id):
    return "fail, Invalid user_id: {}. user_id 长度必须小于64，且必须只含有字母，数字和下划线且字母开头".format(user_id)


def isURL(string):
    result = urlparse(string)
    return result.scheme != '' and result.netloc != ''


def format_source_documents(ori_source_documents):
    source_documents = []
    for inum, doc in enumerate(ori_source_documents):
        source_info = {'file_id': doc.metadata.get('file_id', ''),
                       'file_name': doc.metadata.get('file_name', ''),
                       'content': doc.page_content,
                       'retrieval_query': doc.metadata.get('retrieval_query', ''),
                       # 'kernel': doc.metadata['kernel'],
                       'file_url': doc.metadata.get('file_url', ''),
                       'score': str(doc.metadata['score']),
                       'embed_version': doc.metadata.get('embed_version', ''),
                       'nos_keys': doc.metadata.get('nos_keys', ''),
                       'doc_id': doc.metadata.get('doc_id', ''),
                       'retrieval_source': doc.metadata.get('retrieval_source', ''),
                       'headers': doc.metadata.get('headers', {}),
                       'page_id': doc.metadata.get('page_id', 0),
                       }
        source_documents.append(source_info)
    return source_documents


def format_time_record(time_record):
    token_usage = {}
    time_usage = {}
    for k, v in time_record.items():
        if 'tokens' in k:
            token_usage[k] = v
        else:
            time_usage[k] = v
    if 'rewrite_prompt_tokens' in token_usage:
        if 'prompt_tokens' in token_usage:
            token_usage['prompt_tokens'] += token_usage['rewrite_prompt_tokens']
        if 'total_tokens' in token_usage:
            token_usage['total_tokens'] += token_usage['rewrite_prompt_tokens']
    if 'rewrite_completion_tokens' in token_usage:
        if 'completion_tokens' in token_usage:
            token_usage['completion_tokens'] += token_usage['rewrite_completion_tokens']
        if 'total_tokens' in token_usage:
            token_usage['total_tokens'] += token_usage['rewrite_completion_tokens']
    return {"time_usage": time_usage, "token_usage": token_usage}


def safe_get(req: Request, attr: str, default=None):
    try:
        if attr in req.form:
            return req.form.getlist(attr)[0]
        if attr in req.args:
            return req.args[attr]
        if attr in req.json:
            return req.json[attr]
        # if value := req.form.get(attr):
        #     return value
        # if value := req.args.get(attr):
        #     return value
        # """req.json执行时不校验content-type，body字段可能不能被正确解析为json"""
        # if value := req.json.get(attr):
        #     return value
    except BadRequest:
        logging.warning(f"missing {attr} in request")
    except Exception as e:
        logging.warning(f"get {attr} from request failed:")
        logging.warning(traceback.format_exc())
    return default


def truncate_filename(filename, max_length=200):
    # 获取文件名后缀
    file_ext = os.path.splitext(filename)[1]

    # 获取不带后缀的文件名
    file_name_no_ext = os.path.splitext(filename)[0]

    # 计算文件名长度，注意中文字符
    filename_length = len(filename.encode('utf-8'))

    # 如果文件名长度超过最大长度限制
    if filename_length > max_length:
        debug_logger.warning("文件名长度超过最大长度限制，将截取文件名")
        # 生成一个时间戳标记
        timestamp = str(int(time.time()))
        # 截取文件名
        while filename_length > max_length:
            file_name_no_ext = file_name_no_ext[:-4]
            new_filename = file_name_no_ext + "_" + timestamp + file_ext
            filename_length = len(new_filename.encode('utf-8'))
    else:
        new_filename = filename

    return new_filename


# 同步执行环境下的耗时统计装饰器
def get_time(func):
    def get_time_inner(*arg, **kwargs):
        s_time = time.time()
        res = func(*arg, **kwargs)
        e_time = time.time()
        if 'embed' in func.__name__:
            embed_logger.info('函数 {} 执行耗时: {:.2f} 秒'.format(func.__name__, e_time - s_time))
        elif 'rerank' in func.__name__:
            rerank_logger.info('函数 {} 执行耗时: {:.2f} 秒'.format(func.__name__, e_time - s_time))
        else:
            debug_logger.info('函数 {} 执行耗时: {:.2f} 毫秒'.format(func.__name__, (e_time - s_time) * 1000))
        return res

    return get_time_inner


# 异步执行环境下的耗时统计装饰器
def get_time_async(func):
    @wraps(func)
    async def get_time_async_inner(*args, **kwargs):
        s_time = time.perf_counter()
        res = await func(*args, **kwargs)  # 注意这里使用 await 来调用异步函数
        e_time = time.perf_counter()
        if 'embed' in func.__name__:
            embed_logger.info('函数 {} 执行耗时: {:.2f} 秒'.format(func.__name__, e_time - s_time))
        elif 'rerank' in func.__name__:
            rerank_logger.info('函数 {} 执行耗时: {:.2f} 秒'.format(func.__name__, e_time - s_time))
        else:
            debug_logger.info('函数 {} 执行耗时: {:.2f} 毫秒'.format(func.__name__, (e_time - s_time) * 1000))
        return res

    return get_time_async_inner


def read_files_with_extensions():
    # 获取当前脚本文件的路径
    current_file = os.path.abspath(__file__)

    # 获取当前脚本文件所在的目录
    current_dir = os.path.dirname(current_file)

    # 获取项目根目录
    project_dir = os.path.dirname(os.path.dirname(current_dir))

    directory = project_dir + '/data'

    extensions = ['.md', '.txt', '.pdf', '.jpg', '.docx', '.xlsx', '.eml', '.csv']
    for root, dirs, files in os.walk(directory):
        for file in files:
            if file.endswith(tuple(extensions)):
                file_path = os.path.join(root, file)
                yield file_path


def validate_user_id(user_id):
    if len(user_id) > 64:
        return False
    # 定义正则表达式模式
    pattern = r'^[A-Za-z][A-Za-z0-9_]*$'
    # 检查是否匹配
    if isinstance(user_id, str) and re.match(pattern, user_id):
        return True
    else:
        return False


def num_tokens(text: str, model: str = 'gpt-3.5-turbo-0613') -> int:
    """Return the number of tokens in a string."""
    encoding = tiktoken.encoding_for_model(model)
    return len(encoding.encode(text, disallowed_special=()))


embedding_tokenizer = AutoTokenizer.from_pretrained(LOCAL_EMBED_PATH, local_files_only=True)
rerank_tokenizer = AutoTokenizer.from_pretrained(LOCAL_RERANK_PATH, local_files_only=True)


def num_tokens_embed(text: str) -> int:
    """Return the number of tokens in a string."""
    return len(embedding_tokenizer.encode(text, add_special_tokens=True))


def num_tokens_rerank(text: str) -> int:
    """Return the number of tokens in a string."""
    return len(rerank_tokenizer.encode(text, add_special_tokens=True))


def shorten_data(data):
    # copy data，不要修改原始数据
    data = data.copy()
    try:
        for k, v in data.items():
            if len(str(v)) > 100:
                data[k] = str(v)[:100] + '...'
    except Exception as e:
        debug_logger.error('shorten_data error:', traceback.format_exc())
    return data


def cur_func_name():
    return inspect.currentframe().f_back.f_code.co_name


def num_tokens_from_messages(message_texts, model="gpt-3.5-turbo-0301"):
    encoding = tiktoken.encoding_for_model(model)
    num_tokens = 0
    for message in message_texts:
        # num_tokens += 4  # every message follows <im_start>{role/name}\n{content}<im_end>\n
        # for key, value in message.items():
        num_tokens += len(encoding.encode(message, disallowed_special=()))
        # if key == "name":  # if there's a name, the role is omitted
        # num_tokens += -1  # role is always required and always 1 token
    # num_tokens += 2  # every reply is primed with <im_start>assistant
    return num_tokens


def sent_tokenize(x):
    #  sents_temp = re.split('(：|:|,|，|。|！|\!|\.|？|\?)', x)
    sents_temp = re.split('(。|！|\!|\.|？|\?)', x)
    sents = []
    for i in range(len(sents_temp) // 2):
        sent = sents_temp[2 * i] + sents_temp[2 * i + 1]
        sents.append(sent)
    return sents


def clear_string(str):
    # 只保留中文、英文、数字
    str = re.sub(r"[^\u4e00-\u9fa5a-zA-Z0-9]", "", str)
    return str


def simplify_filename(filename, max_length=40):
    if len(filename) <= max_length:
        # 如果文件名长度小于等于最大长度，直接返回原文件名
        return filename

    # 分离文件的基本名和扩展名
    name, extension = filename.rsplit('.', 1)
    extension = '.' + extension  # 将点添加回扩展名

    # 计算头部和尾部的保留长度
    part_length = (max_length - len(extension) - 1) // 2  # 减去扩展名长度和破折号的长度
    end_start = -part_length if part_length else None

    # 构建新的简化文件名
    simplified_name = f"{name[:part_length]}-{name[end_start:]}" if part_length else name[:max_length - 1]

    return f"{simplified_name}{extension}"


# 对比两个字符串，只保留字母数字和中文，返回是否一致
def clear_string_is_equal(str1, str2):
    str1 = clear_string(str1)
    str2 = clear_string(str2)
    return str1 == str2


def correct_kb_id(kb_id):
    if not kb_id:
        return kb_id
    # 如果kb_id末尾不是KB_SUFFIX,则加上
    if KB_SUFFIX not in kb_id:
        if kb_id.endswith('_FAQ'):  # KBc86eaa3f278f4ef9908780e8e558c6eb_FAQ
            return kb_id.split('_FAQ')[0] + KB_SUFFIX + '_FAQ'
        else:  # KBc86eaa3f278f4ef9908780e8e558c6eb
            return kb_id + KB_SUFFIX
    else:
        return kb_id


def clear_kb_id(kb_id):
    return kb_id.replace(KB_SUFFIX, '')


def string_bytes_length(string):
    return len(string.encode('utf-8'))


def export_qalogs_to_excel(qalogs, columns, filename: str):
    # 将查询结果转换为 DataFrame
    df = pd.DataFrame(qalogs, columns=columns)

    # 写入 Excel 文件
    root_path = os.path.dirname(UPLOAD_ROOT_PATH) + '/saved_qalogs'
    if not os.path.exists(root_path):
        os.makedirs(root_path)

    file_path = os.path.join(root_path, filename)
    df.to_excel(file_path, index=False)

    # 使用 openpyxl 调整列宽
    workbook = load_workbook(filename=file_path)
    worksheet = workbook.active

    for column_cells in worksheet.columns:
        length = max(len(str(cell.value)) for cell in column_cells)
        worksheet.column_dimensions[get_column_letter(column_cells[0].column)].width = length

    workbook.save(file_path)
    debug_logger.info(f"Data exported to {file_path} successfully.")
    return file_path


def check_user_id_and_user_info(user_id, user_info):
    if user_id is None or user_info is None:
        msg = "fail, user_id 或 user_info 为 None"
        return False, msg
    if not validate_user_id(user_id):
        msg = get_invalid_user_id_msg(user_id)
        return False, msg
    if not user_info.isdigit():
        msg = "fail, user_info 必须是纯数字"
        return False, msg
    return True, 'success'


def cosine_similarity(embedding1, embedding2):
    embedding1 = np.array(embedding1)
    embedding2 = np.array(embedding2)
    # 计算两个向量的点积
    dot_product = np.dot(embedding1, embedding2)
    # 计算两个向量的模
    norm_embedding1 = np.linalg.norm(embedding1)
    norm_embedding2 = np.linalg.norm(embedding2)
    # 计算余弦相似度
    similarity = dot_product / (norm_embedding1 * norm_embedding2)
    # 将余弦相似度映射到0-1之间
    similarity_mapped = (similarity + 1) / 2
    return similarity_mapped


def get_table_infos(markdown_str):
    lines = markdown_str.split('\n')
    if len(lines) < 2:
        return None
    head_line = None
    end_line = None
    for i in range(len(lines) - 1):
        if '|' in lines[i] and '|' in lines[i + 1]:
            separator_line = lines[i + 1].strip()
            if separator_line.startswith('|') and separator_line.endswith('|'):
                separator_parts = separator_line[1:-1].split('|')
                if all(part.strip().startswith('-') and len(part.strip()) >= 3 for part in separator_parts):
                    head_line = i
                    break
    for i in range(len(lines)):
        if '|' in lines[i]:
            separator_line = lines[i].strip()
            if separator_line.startswith('|') and separator_line.endswith('|'):
                end_line = i
    if head_line is None or end_line is None:
        return None
    return {"head_line": head_line, "end_line": end_line, "head": lines[head_line] + '\n' + lines[head_line + 1],
            "lines": lines}


def get_time_range(time_start=None, time_end=None, default_days=7):
    """
    获取时间范围。如果给定的时间范围不完整，将使用默认值（最近7天）。

    :param time_start: 起始时间，格式为 "YYYY-MM-DD" 或 "YYYY-MM-DD HH:MM:SS"
    :param time_end: 结束时间，格式为 "YYYY-MM-DD" 或 "YYYY-MM-DD HH:MM:SS"
    :param default_days: 如果未提供时间范围，默认的天数范围
    :return: 包含起始时间和结束时间的元组，格式为 ("YYYY-MM-DD HH:MM:SS", "YYYY-MM-DD HH:MM:SS")
    """

    date_pattern = r'^\d{4}-\d{2}-\d{2}$'
    now = datetime.now()

    # 验证 time_start 格式
    if time_start:
        if not re.match(date_pattern, time_start):
            return None
        if len(time_start) == 10:
            time_start = time_start + " 00:00:00"
    else:
        time_start = (now - timedelta(days=default_days)).strftime("%Y-%m-%d 00:00:00")

    # 验证 time_end 格式
    if time_end:
        if not re.match(date_pattern, time_end):
            return None
        if len(time_end) == 10:
            time_end = time_end + " 23:59:59"
    else:
        time_end = now.strftime("%Y-%m-%d 23:59:59")

    return (time_start, time_end)


def deduplicate_documents(source_docs):
    unique_docs = set()
    deduplicated_docs = []
    for doc in source_docs:
        if doc.page_content not in unique_docs:
            unique_docs.add(doc.page_content)
            deduplicated_docs.append(doc)
    return deduplicated_docs


def get_all_subpages(url):
    headers = {
        'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
    }

    try:
        response = requests.get(url, headers=headers, timeout=10)
        response.raise_for_status()
    except requests.RequestException as e:
        print(f"Error fetching {url}: {e}")
        return []

    soup = BeautifulSoup(response.text, 'html.parser')
    links = soup.find_all('a', href=True)
    subpages = set()

    for link in links:
        href = link['href']
        full_url = urljoin(url, href)
        subpages.add(full_url)

    return list(subpages)


def html_to_markdown(html_content):
    # 创建HTML到文本转换器
    h = html2text.HTML2Text()

    # 配置转换器
    h.ignore_images = True
    h.ignore_emphasis = True
    h.ignore_links = True
    h.body_width = 0  # 禁用换行
    h.tables = True  # 保留表格

    # 转换HTML到Markdown
    markdown = h.handle(html_content)

    # 删除所有图片标记
    markdown = re.sub(r'!\[.*?\]\(.*?\)', '', markdown)

    # 删除所有链接标记，保留文字
    markdown = re.sub(r'\[([^\]]*)\]\(.*?\)', r'\1', markdown)

    # 删除多余的空行，但保留表格结构
    # markdown = re.sub(r'(\n\s*){3,}', '\n\n', markdown)

    # 删除行首的特殊字符（如*、-等），但保留表格的|符号
    # markdown = re.sub(r'^(?!\|)\s*[-*]\s+', '', markdown, flags=re.MULTILINE)

    return markdown.strip()


def fast_estimate_file_char_count(file_path):
    """
    快速估算文件的字符数，如果超过max_chars则返回False，否则返回True
    """
    file_extension = os.path.splitext(file_path)[1].lower()

    try:
        if file_extension in ['.md', '.txt', '.csv']:
            with open(file_path, 'rb') as file:
                raw = file.read(1024)
                encoding = chardet.detect(raw)['encoding']
            with open(file_path, 'r', encoding=encoding) as file:
                char_count = sum(len(line) for line in file)

        elif file_extension == '.pdf':
            doc = fitz.open(file_path)
            char_count = sum(len(page.get_text()) for page in doc)
            doc.close()

        elif file_extension in ['.jpg', '.png', '.jpeg']:
            # 图片文件无法准确估算字符数，返回True让后续OCR处理
            return True

        elif file_extension == '.docx':
            text = docx2txt.process(file_path)
            char_count = len(text)

        elif file_extension == '.xlsx':
            wb = openpyxl.load_workbook(file_path, read_only=True)
            char_count = sum(len(str(cell.value or '')) for sheet in wb for row in sheet.iter_rows() for cell in row)
            wb.close()

        elif file_extension == '.pptx':
            prs = Presentation(file_path)
            char_count = sum(
                len(shape.text) for slide in prs.slides for shape in slide.shapes if hasattr(shape, 'text'))

        elif file_extension == '.eml':
            with open(file_path, 'r', encoding='utf-8') as file:
                msg = email.message_from_file(file)
                char_count = len(str(msg))

        else:
            # 不支持的文件类型
            return False

        return char_count

    except Exception as e:
        print(f"Error processing file {file_path}: {str(e)}")
        return None
