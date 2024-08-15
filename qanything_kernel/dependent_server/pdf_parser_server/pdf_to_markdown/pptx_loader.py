from io import BytesIO
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from rapidocr_onnxruntime import RapidOCR
import numpy as np
import cv2
from PIL import Image
from qanything_kernel.utils.loader.pdf_to_markdown.filetype import html2markdown


def pptx_table_to_html(table):
    html_table = '<table>\n'
    for row in table.rows:
        html_table += '  <tr>\n'
        for cell in row.cells:
            cell_text = ''.join(text_frame.text for text_frame in cell.text_frame.paragraphs)
            html_table += '    <td>{}</td>\n'.format(cell_text)
        html_table += '  </tr>\n'
    html_table += '</table>\n'
    return html_table


class PPTParser(object):
    def __init__(self):
        super().__init__()
        self.ocr_rapid = RapidOCR()

    def __extract(self, shape):
        if shape.shape_type == 19:
            tb = shape.table
            tb_html = pptx_table_to_html(tb)
            new_html = tb_html

            tb_markdown = html2markdown(new_html)

            return tb_markdown

        if shape.has_text_frame:
            return shape.text_frame.text

        if shape.shape_type == 6:
            texts = []
            for p in sorted(shape.shapes, key=lambda x: (x.top // 10, x.left)):
                t = self.__extract(p)
                if t:
                    texts.append(t)
            return "\n".join(texts)

    def __extract_img(self, slide, shape, array_lst):
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image_data = shape.image.blob
            image = Image.open(BytesIO(image_data))
            array = cv2.cvtColor(np.asarray(image), cv2.COLOR_RGB2BGR)
            # cv2.imwrite(f'./tmp/{slide.slide_id}_{shape.shape_id}.png',array)
            array_lst.append(array)
            return array_lst, True
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for s in shape.shapes:
                self.__extract_img(slide, s, array_lst)
            return array_lst, True
        else:
            return [], False

    def __call__(self, fnm, from_page, to_page, callback=None):
        ppt = Presentation(fnm) if isinstance(
            fnm, str) else Presentation(
            BytesIO(fnm))
        txts = []
        self.total_page = len(ppt.slides)
        # print(self.total_page)
        for i, slide in enumerate(ppt.slides):
            if i < from_page:
                continue
            if i >= to_page:
                break
            texts = []
            for shape in sorted(
                    slide.shapes, key=lambda x: (x.top // 10, x.left)):
                txt = self.__extract(shape)
                array_lst = []
                image_lst, has_image = self.__extract_img(slide, shape, array_lst)
                if has_image:
                    ocr_res_lst = []
                    for img in image_lst:
                        ocr_res = self.__ocr(img)
                        if ocr_res != '':
                            ocr_res_lst.append(ocr_res)
                    if len(ocr_res_lst) != 0:
                        ocr_res = '\n'.join(ocr_res_lst)
                    else:
                        ocr_res = None
                else:
                    ocr_res = None
                if txt:
                    texts.append(txt)
                if ocr_res:
                    texts.append(ocr_res)
            txts.append("\n".join(texts))
        return txts

    def __ocr(self, image):
        bxs_rapid = self.ocr_rapid(image)
        text = ''
        if bxs_rapid[0] == None:
            return ''
        for box in bxs_rapid[0]:
            text += (box[1] + '\n')
        return text


if __name__ == '__main__':
    path = '第一次培训_2024年第一次培训大模型业务介绍V1.0.pptx'
    parser = PPTParser()
    res = parser(fnm=path, from_page=12, to_page=13)
    # print(res)
    for text in res:
        print('********************')
        print(text)
